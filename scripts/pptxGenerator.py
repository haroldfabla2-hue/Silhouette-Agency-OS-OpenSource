#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Autonomous PowerPoint Slide Deck Generator for Silhouette OS
Alberto Farah Agency
"""

import sys
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation(data_path, out_path):
    if not os.path.exists(data_path):
        print(f"Error: input file {data_path} not found.")
        sys.exit(1)
        
    with open(data_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
        
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 widescreen default
    prs.slide_height = Inches(7.5)
    
    # Brand Colors (Montserrat & Inter Theme)
    NEGRO = RGBColor(10, 10, 10)
    BLANCO = RGBColor(245, 245, 245)
    VERDE_LIMA = RGBColor(57, 255, 20)
    GRIS_OSCURO = RGBColor(30, 30, 30)
    GRIS_LIGHT = RGBColor(180, 180, 180)
    
    slides_data = data.get("slides", [])
    
    for idx, slide_info in enumerate(slides_data):
        layout_name = slide_info.get("layout", "default")
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        
        # Apply dark background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = NEGRO
        bg.line.fill.background()
        
        title_text = slide_info.get("title", "").upper()
        content_text = slide_info.get("content", "")
        
        if idx == 0 or layout_name == "title":
            # Title slide layout
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(2))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = 'Montserrat'
            p.font.size = Pt(60)
            p.font.color.rgb = BLANCO
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Accent Line
            line = slide.shapes.add_connector(
                1, Inches(4.5), Inches(4.5), Inches(8.83), Inches(4.5)
            )
            line.line.color.rgb = VERDE_LIMA
            line.line.width = Pt(5)
            
            # Subtitle
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.33), Inches(1.5))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = slide_info.get("subtitle", slide_info.get("content", ""))
            p_sub.font.name = 'Inter'
            p_sub.font.size = Pt(22)
            p_sub.font.color.rgb = GRIS_LIGHT
            p_sub.alignment = PP_ALIGN.CENTER
            
        elif layout_name == "two-column":
            # Header
            header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
            tf_head = header_box.text_frame
            p_head = tf_head.paragraphs[0]
            p_head.text = title_text
            p_head.font.name = 'Montserrat'
            p_head.font.size = Pt(36)
            p_head.font.color.rgb = VERDE_LIMA
            p_head.font.bold = True
            
            # Two Columns
            cols = content_text.split("---")
            left_text = cols[0].strip() if len(cols) > 0 else ""
            right_text = cols[1].strip() if len(cols) > 1 else ""
            
            # Left Box
            l_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
            tf_l = l_box.text_frame
            tf_l.word_wrap = True
            p_l = tf_l.paragraphs[0]
            p_l.text = left_text
            p_l.font.name = 'Inter'
            p_l.font.size = Pt(18)
            p_l.font.color.rgb = BLANCO
            
            # Right Box
            r_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.5))
            tf_r = r_box.text_frame
            tf_r.word_wrap = True
            p_r = tf_r.paragraphs[0]
            p_r.text = right_text
            p_r.font.name = 'Inter'
            p_r.font.size = Pt(18)
            p_r.font.color.rgb = BLANCO
            
        elif layout_name == "quote":
            # Giant quotation mark
            quote_mark = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(11.33), Inches(1))
            tf_q = quote_mark.text_frame
            p_q = tf_q.paragraphs[0]
            p_q.text = "“"
            p_q.font.name = 'Montserrat'
            p_q.font.size = Pt(120)
            p_q.font.color.rgb = VERDE_LIMA
            p_q.alignment = PP_ALIGN.CENTER
            
            # Quote Text
            quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.33), Inches(3.5))
            tf_quote = quote_box.text_frame
            tf_quote.word_wrap = True
            p_quote = tf_quote.paragraphs[0]
            p_quote.text = content_text
            p_quote.font.name = 'Inter'
            p_quote.font.size = Pt(28)
            p_quote.font.color.rgb = BLANCO
            p_quote.font.italic = True
            p_quote.alignment = PP_ALIGN.CENTER
            
            # Author / Cite
            if title_text and title_text != "SLIDE":
                cite_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(10.33), Inches(1))
                tf_cite = cite_box.text_frame
                p_cite = tf_cite.paragraphs[0]
                p_cite.text = f"— {title_text}"
                p_cite.font.name = 'Montserrat'
                p_cite.font.size = Pt(20)
                p_cite.font.color.rgb = VERDE_LIMA
                p_cite.font.bold = True
                p_cite.alignment = PP_ALIGN.CENTER
                
        elif layout_name == "list":
            # Header
            header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
            tf_head = header_box.text_frame
            p_head = tf_head.paragraphs[0]
            p_head.text = title_text
            p_head.font.name = 'Montserrat'
            p_head.font.size = Pt(36)
            p_head.font.color.rgb = BLANCO
            p_head.font.bold = True
            
            # Bullet List
            list_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
            tf_list = list_box.text_frame
            tf_list.word_wrap = True
            
            lines = [line.strip() for line in content_text.split("\n") if line.strip()]
            for idx, line_val in enumerate(lines):
                if idx == 0:
                    p_item = tf_list.paragraphs[0]
                else:
                    p_item = tf_list.add_paragraph()
                p_item.text = line_val.replace("•", "").strip()
                p_item.font.name = 'Inter'
                p_item.font.size = Pt(20)
                p_item.font.color.rgb = BLANCO
                p_item.space_after = Pt(14)
                p_item.level = 0
                
        else:
            # Standard Text Slide Layout
            header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
            tf_head = header_box.text_frame
            p_head = tf_head.paragraphs[0]
            p_head.text = title_text
            p_head.font.name = 'Montserrat'
            p_head.font.size = Pt(36)
            p_head.font.color.rgb = VERDE_LIMA
            p_head.font.bold = True
            
            body_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
            tf_body = body_box.text_frame
            tf_body.word_wrap = True
            p_body = tf_body.paragraphs[0]
            p_body.text = content_text
            p_body.font.name = 'Inter'
            p_body.font.size = Pt(20)
            p_body.font.color.rgb = BLANCO
            
        # Optional Slide Image inclusion if provided and exists
        image_path = slide_info.get("imagePath", "")
        if image_path and os.path.exists(image_path):
            try:
                # Add image in the right column space if layout isn't title
                if layout_name not in ["title", "quote"]:
                    slide.shapes.add_picture(image_path, Inches(7.5), Inches(2.0), width=Inches(5.0))
            except Exception as e:
                print(f"Warning: could not add image {image_path}: {e}")

    # Save output PPTX
    prs.save(out_path)
    print(f"Success: Presentation generated at {out_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python pptxGenerator.py <input_json_path> <output_pptx_path>")
        sys.exit(1)
    build_presentation(sys.argv[1], sys.argv[2])
